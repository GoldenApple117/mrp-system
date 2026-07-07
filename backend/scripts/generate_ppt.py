"""生成 MRP II 系统周工作总结 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# 颜色方案
C_PRIMARY = RGBColor(0x25, 0x63, 0xEB)   # 蓝色
C_DARK = RGBColor(0x1A, 0x23, 0x32)       # 深色
C_GRAY = RGBColor(0x6B, 0x7C, 0x93)       # 灰色
C_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)      # 浅灰背景
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN = RGBColor(0x05, 0x96, 0x69)
C_AMBER = RGBColor(0xD9, 0x77, 0x06)
C_VIOLET = RGBColor(0x7C, 0x3A, 0xED)
C_ACCENT = RGBColor(0x08, 0x91, 0xB2)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=C_DARK, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, items, font_size=13, color=C_DARK, spacing=Pt(6), prefix=''):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{prefix}{item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = spacing
    return txBox

def add_kpi_card(slide, left, top, num, label, num_color=C_PRIMARY):
    card = add_rect(slide, left, top, Inches(1.9), Inches(1.1), C_WHITE)
    card.shadow.inherit = False
    # number
    add_textbox(slide, left + Inches(0.15), top + Inches(0.1), Inches(1.6), Inches(0.55),
                str(num), font_size=24, bold=True, color=num_color, alignment=PP_ALIGN.CENTER)
    # label
    add_textbox(slide, left + Inches(0.15), top + Inches(0.55), Inches(1.6), Inches(0.4),
                label, font_size=9, color=C_GRAY, alignment=PP_ALIGN.CENTER)

def add_phase_block(slide, left, top, title, desc, color=C_ACCENT):
    card = add_rect(slide, left, top, Inches(3.9), Inches(1.0), C_WHITE)
    card.shadow.inherit = False
    # color bar on left
    add_rect(slide, left, top, Inches(0.08), Inches(1.0), color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.08), Inches(3.5), Inches(0.3),
                title, font_size=12, bold=True, color=C_DARK)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.35), Inches(3.5), Inches(0.6),
                desc, font_size=10, color=C_GRAY)

# ===========================================================
# Slide 1: 封面
# ===========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, C_PRIMARY)
# Title
add_textbox(slide, Inches(1), Inches(1.8), Inches(8), Inches(0.8),
            "MRP II 物料需求计划系统", font_size=36, bold=True, color=C_WHITE)
# Subtitle
add_textbox(slide, Inches(1), Inches(2.7), Inches(8), Inches(0.6),
            "周工作总结汇报", font_size=24, bold=False, color=C_WHITE)
# Divider line
add_rect(slide, Inches(1), Inches(3.6), Inches(2), Inches(0.04), C_WHITE)
# Info
add_textbox(slide, Inches(1), Inches(3.9), Inches(8), Inches(0.4),
            "汇报人：金崧（算法小组 · P1 实习生）", font_size=14, color=C_WHITE)
add_textbox(slide, Inches(1), Inches(4.4), Inches(8), Inches(0.4),
            "时间周期：2026年7月1日 ~ 7月7日 · 入职第13~19天", font_size=14, color=C_WHITE)

# ===========================================================
# Slide 2: 目录
# ===========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_textbox(slide, Inches(1), Inches(1.0), Inches(6), Inches(0.6),
            "目  录", font_size=28, bold=True, color=C_DARK)
add_rect(slide, Inches(1), Inches(1.7), Inches(1.5), Inches(0.04), C_PRIMARY)

items = [
    ("01", "系统当前数据概览"),
    ("02", "开发历程回顾（七日迭代）"),
    ("03", "关键技术决策"),
    ("04", "个人心得与当前状态"),
]
for i, (num, title) in enumerate(items):
    y = Inches(2.4) + Inches(0.8) * i
    add_rect(slide, Inches(1.5), y, Inches(0.5), Inches(0.5), C_PRIMARY)
    add_textbox(slide, Inches(1.5), y + Inches(0.05), Inches(0.5), Inches(0.4),
                num, font_size=14, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.3), y + Inches(0.05), Inches(6), Inches(0.4),
                title, font_size=16, color=C_DARK)

# ===========================================================
# Slide 3: 系统数据概览
# ===========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
# Section header
add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), C_PRIMARY)
add_textbox(slide, Inches(1), Inches(0.2), Inches(8), Inches(0.6),
            "01  系统当前数据概览", font_size=24, bold=True, color=C_WHITE)
add_textbox(slide, Inches(1), Inches(0.7), Inches(8), Inches(0.4),
            "系统规模与质量指标", font_size=13, color=C_WHITE)

# KPI cards
kpis = [
    ("170+", "API 端点数", C_PRIMARY),
    ("27", "数据库表", C_GREEN),
    ("973", "物料数", C_AMBER),
    ("858", "BOM 行数", C_VIOLET),
    ("1332", "采购订单", C_ACCENT),
    ("97.1%", "E2E 测试通过率", C_GREEN),
]
for i, (num, label, color) in enumerate(kpis):
    col = i % 3
    row = i // 3
    x = Inches(1) + Inches(3.8) * col
    y = Inches(2.0) + Inches(1.8) * row
    add_kpi_card(slide, x, y, num, label, color)

# Coverage text
add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
            "测试覆盖：登录 → 物料 / BOM / 库存 / ABC → 销售订单审核自动MPS → MRP运算 → IQC/NCR → 工单执行 → 批次追溯 / 序列号 → 看板 / OEE / 安灯 → 设备 / 模具 → 报表",
            font_size=11, color=C_GRAY)

# ===========================================================
# Slide 4: 开发历程章节页
# ===========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_ACCENT)
add_textbox(slide, Inches(1), Inches(2.5), Inches(10), Inches(0.8),
            "02  开发历程回顾", font_size=32, bold=True, color=C_WHITE)
add_textbox(slide, Inches(1), Inches(3.5), Inches(10), Inches(0.5),
            "七日迭代与功能演进", font_size=16, color=C_WHITE)

# ===========================================================
# Slides 5-10: 每日进展
# ===========================================================

days = [
    ("7月1日 · 周三", "Phase 2 核心模块联调", [
        "销售订单管理：客户/订单模型、自动计算总应收、出货/收款双状态跟踪",
        "费用合计：项目→模块→零件三层展示、206个零件、总金额¥90,750",
        "前端重构：三级折叠视图（采购/物料/库存）、div+v-show替代el-collapse",
        "数据联动：产品出库→出货状态、客户付款→收款状态、销售订单→MPS推送",
    ], C_PRIMARY),
    ("7月2日 · 周四", "系统定位审视 — 暂停开发先想清楚", [
        "分析三种使用路径：全员分权限 / 集中填报专人操作 / Excel导入备份",
        "核心原则确立：系统应让用户少做一步，而非多做一步",
        "与周圆圆前辈交流公司WPS云文档使用方式",
        "向领导提出三个问题：使用人群、使用方式、核心价值",
    ], C_VIOLET),
    ("7月3日 · 周五", "v1.7安全加固 & 云端部署", [
        "JWT认证系统：Canvas粒子登录页、bcrypt密码验证、24h Token、路由守卫",
        "双角色权限管理：管理员/普通用户、前端遮罩+后端403双重保护",
        "P0安全加固：JWT Secret环境变量化、smtp_config进.gitignore",
        "Railway云端部署、定时MRP（每天06:00自动运算）、向王琛前辈提交操作指南",
    ], C_GREEN),
    ("7月4日 · 周六", "云端UAT验收 & 产品方向明确", [
        "25项自动化验收，通过率92%（2项未通过实为正确的业务保护逻辑）",
        "Excel导出500修复：中文文件名Linux容器latin-1编码→RFC 5987方案",
        "Docker缓存问题：修改RUN行注释强制重建",
        "王琛前辈明确方向：四计划协同、计划滚动+备货制、云文档后期停用",
    ], C_AMBER),
    ("7月6日 · 周一", "业务闭环构建 — 15个新API", [
        "生产修复：MPS创建500 + MRP运算500（JSON溢出135KB→LONGTEXT）",
        "新增15个API：销售审批流、采购收货双流水、生产领退料、工单报工、成本核算",
        "分析功能：ABC帕累托分析、供应商绩效、RCCP产能分析、ATP可承诺量",
        "环境迁移：SQLite → MySQL 8.4，实现本地与云端数据库一致",
    ], C_ACCENT),
]

for day_title, main_title, items, color in days:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_WHITE)
    # Left color bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), color)
    # Date tag
    add_rect(slide, Inches(1), Inches(0.8), Inches(2.2), Inches(0.45), color)
    add_textbox(slide, Inches(1), Inches(0.83), Inches(2.2), Inches(0.4),
                day_title, font_size=11, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
    # Main title
    add_textbox(slide, Inches(1), Inches(1.5), Inches(10), Inches(0.5),
                main_title, font_size=22, bold=True, color=C_DARK)
    # Divider
    add_rect(slide, Inches(1), Inches(2.2), Inches(2), Inches(0.03), color)
    # Items
    add_multi_text(slide, Inches(1), Inches(2.6), Inches(11), Inches(4.5), items, font_size=14, color=C_DARK, prefix='▪ ')

# ===========================================================
# Slide: 7月7日 Part 1 (Phase 6-8)
# ===========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), C_PRIMARY)
add_rect(slide, Inches(1), Inches(0.8), Inches(2.2), Inches(0.45), C_PRIMARY)
add_textbox(slide, Inches(1), Inches(0.83), Inches(2.2), Inches(0.4),
            "7月7日 · 周二", font_size=11, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(1.5), Inches(10), Inches(0.5),
            "制造执行深化 — Phase 6~10（17页高能）", font_size=22, bold=True, color=C_DARK)
add_rect(slide, Inches(1), Inches(2.2), Inches(2), Inches(0.03), C_PRIMARY)

# Phase blocks in 2x3 grid
phases_data = [
    ("Phase 6", "工序级工单执行", "WorkOrderOperation模型（19字段）\n6个操作端点：初始化/开工/报工/完工/跳过\n工艺路线自动展开、逐序执行、下一道自动解锁", C_PRIMARY),
    ("Phase 7", "质检体系", "IQC来料检 / PQC过程检 / OQC出货检\n不合格自动生成NCR记录\n修复采购收货+IQC重复累加库存bug", C_GREEN),
    ("Phase 8", "批次追溯", "批次记录与序列号管理\n正向追溯：批次→下游客户\n反向追溯：序列号→上游供应商", C_AMBER),
    ("Phase 9", "生产看板", "车间总览KPI卡片（活跃工单/进行中）\nOEE设备综合效率计算\n安灯异常告警系统", C_VIOLET),
    ("Phase 10", "设备模具管理", "设备台账（型号/状态/使用部门）\n模具台账（寿命/使用次数）\n保养计划（周期/到期提醒）", C_ACCENT),
    ("E2E测试", "全流程回归测试", "34项API测试，33通过/0失败\n通过率97.1%\n测试脚本：backend/scripts/e2e_test.py", C_DARK),
]
for i, (tag, title, desc, color) in enumerate(phases_data):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6.2) * col
    y = Inches(2.6) + Inches(1.55) * row
    # Tag badge
    add_rect(slide, x, y, Inches(1.2), Inches(0.35), color)
    add_textbox(slide, x, y + Inches(0.02), Inches(1.2), Inches(0.3),
                tag, font_size=10, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(1.4), y, Inches(4.5), Inches(0.35),
                title, font_size=14, bold=True, color=C_DARK)
    # Desc
    add_textbox(slide, x + Inches(1.4), y + Inches(0.4), Inches(4.5), Inches(1.0),
                desc, font_size=10.5, color=C_GRAY)

# ===========================================================
# Slide: 7月7日 Part 2 (前端+E2E)
# ===========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), C_PRIMARY)
add_textbox(slide, Inches(1), Inches(0.2), Inches(10), Inches(0.6),
            "前端优化 & E2E 测试", font_size=22, bold=True, color=C_WHITE)
add_textbox(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.4),
            "7月7日 · 周二", font_size=13, color=C_WHITE)

# Frontend
add_textbox(slide, Inches(1), Inches(1.6), Inches(5), Inches(0.4),
            "前端优化（精简合并）", font_size=18, bold=True, color=C_DARK)
add_rect(slide, Inches(1), Inches(2.1), Inches(1.2), Inches(0.03), C_ACCENT)
fe_items = [
    "BOM导入弹窗：4个Tab→单Excel上传（删除241行废弃代码）",
    "页面合并：Dashboard + ShopFloor → 统一\"生产驾驶舱\"",
    "菜单改名：\"报表分析\"→\"报表中心\"，与页面标题保持一致",
]
add_multi_text(slide, Inches(1), Inches(2.3), Inches(5), Inches(2.5), fe_items, font_size=13, prefix='▪ ')

# E2E test
add_textbox(slide, Inches(6.5), Inches(1.6), Inches(5.5), Inches(0.4),
            "E2E 全流程回归测试", font_size=18, bold=True, color=C_DARK)
add_rect(slide, Inches(6.5), Inches(2.1), Inches(1.2), Inches(0.03), C_GREEN)

# Result cards
for i, (label, val, color) in enumerate([
    ("通过", "33项", C_GREEN),
    ("失败", "0项", RGBColor(0xD9, 0x2D, 0x20)),
    ("跳过", "1项", C_AMBER),
    ("通过率", "97.1%", C_PRIMARY),
]):
    x = Inches(6.5) + Inches(1.6) * i
    add_kpi_card(slide, x, Inches(2.4), label, val, color)

# Coverage
add_textbox(slide, Inches(6.5), Inches(3.8), Inches(5.5), Inches(2.5),
            "覆盖链路：登录 → 物料CRUD → BOM树 → 库存ABC → 销售订单→MPS → MRP → IQC/NCR → 工单/工序 → 批次追溯 → 序列号 → 车间看板 → OEE → 安灯 → 设备模具 → 报表",
            font_size=11, color=C_GRAY)

# ===========================================================
# Slide: 关键技术决策
# ===========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), C_VIOLET)
add_textbox(slide, Inches(1), Inches(0.2), Inches(10), Inches(0.6),
            "03  关键技术决策", font_size=24, bold=True, color=C_WHITE)
add_textbox(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.4),
            "在开发过程中遇到的四个关键选择", font_size=13, color=C_WHITE)

decisions = [
    ("SQLite → MySQL 迁移",
     "本地环境与云端保持一致，避免\"本地正常、云端500\"",
     "问题：SQLite的LONGTEXT不兼容、DELETE顺序差异、Alembic方言不一致\n解法：init_db()按数据库类型分流，模型用通用类型"),
    ("工单成本：实时计算",
     "每次GET/cost JOIN三张表，避免预聚合的数据不一致",
     "原因：成本在工单执行期间动态变化，预聚合需要每次变动刷新快照\n性能：中小规模工单（几十到几百条）下耗时<10ms"),
    ("采购收货：双流水",
     "合格品写采购入库、不合格品写质检不合格，两笔独立记录",
     "优势：供应商质量合格率可直接按type聚合，不需要解析备注\n后续可关联质检报告，便于批次追溯"),
    ("Excel导出：RFC 5987",
     "中文文件名在Linux容器的编码兼容方案",
     "根因：Starlette默认latin-1编码HTTP响应头\n方案：Content-Disposition用filename*=UTF-8''格式编码"),
]
for i, (title, what, why) in enumerate(decisions):
    x = Inches(0.5) + Inches(3.2) * (i % 2)
    y = Inches(1.8) + Inches(2.7) * (i // 2)
    # Card background
    card = add_rect(slide, x, y, Inches(3.0), Inches(2.5), C_LIGHT)
    # Title
    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.6), Inches(0.35),
                title, font_size=13, bold=True, color=C_VIOLET)
    # What
    add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.6), Inches(0.5),
                what, font_size=10, bold=True, color=C_GREEN)
    # Why
    add_textbox(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.6), Inches(1.3),
                why, font_size=9.5, color=C_GRAY)

# ===========================================================
# Slide: 个人心得
# ===========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), C_VIOLET)
add_textbox(slide, Inches(1), Inches(0.2), Inches(10), Inches(0.6),
            "04  个人心得与反思", font_size=24, bold=True, color=C_WHITE)

reflections = [
    ("先想清楚\"谁来用\"，再决定\"做什么\"",
     "周四那天暂停开发，审视系统定位——如果连使用场景都没想清楚，功能再多也只是自娱自乐。和王琛前辈沟通后明确了方向：不是替代云文档，而是做云文档做不到的事——打通订单/生产/物料/采购四计划的数据协同。"),
    ("开发环境与生产环境保持一致",
     "一周内遇到至少三次\"本地正常、云端500\"：LONGTEXT差异、DELETE顺序、HTTP编码。越早统一开发环境和生产环境，这种调试成本越低。花了半天把SQLite迁移到MySQL 8.4，过程折腾但长痛不如短痛。"),
    ("批量开发后必须跑全量回归",
     "Phase 6~10连续写完只测了新接口，今天写E2E脚本时踩了好几个路径不一致的坑。如果昨天花15分钟跑一遍自动化测试，今天就不用反复改测试脚本。提前暴露跨模块影响才是高效的。"),
    ("真实数据才能暴露边界",
     "MRP运算用TEXT(65KB)存JSON，本地23个用例从未触发上限。生产数据135KB直接报DataError 1406。测试数据的规模和复杂度一定要和真实数据对齐。"),
]
for i, (title, body) in enumerate(reflections):
    y = Inches(1.6) + Inches(1.4) * i
    # Numbered badge
    add_rect(slide, Inches(0.8), y, Inches(0.4), Inches(0.4), C_VIOLET)
    add_textbox(slide, Inches(0.8), y + Inches(0.03), Inches(0.4), Inches(0.35),
                f"0{i+1}", font_size=11, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, Inches(1.5), y - Inches(0.05), Inches(11), Inches(0.35),
                title, font_size=14, bold=True, color=C_DARK)
    # Body
    add_textbox(slide, Inches(1.5), y + Inches(0.3), Inches(11), Inches(0.9),
                body, font_size=10.5, color=C_GRAY)

# ===========================================================
# Slide: 当前状态
# ===========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), C_PRIMARY)
add_textbox(slide, Inches(1), Inches(0.2), Inches(10), Inches(0.6),
            "当前状态", font_size=24, bold=True, color=C_WHITE)

# Done
add_textbox(slide, Inches(1), Inches(1.5), Inches(5), Inches(0.5),
            "✅ 已完成", font_size=20, bold=True, color=C_GREEN)
done_items = [
    "本地业务闭环已打通（销售→MPS→MRP→采购→工单→质检→追溯→看板→报表）",
    "制造执行模块（Phase 6~10）全部完成，新增约25个API，总API ≥ 170",
    "E2E全流程测试通过率97.1%（34项中33项通过）",
    "前端优化：页面合并、弹窗精简、菜单改名",
    "本地开发环境切换为MySQL 8.4，与云端一致",
]
add_multi_text(slide, Inches(1), Inches(2.1), Inches(5.5), Inches(3.5), done_items, font_size=12, prefix='✅ ', color=C_GREEN)

# Todo
add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5), Inches(0.5),
            "⚠️ 待办事项", font_size=20, bold=True, color=C_AMBER)
todo_items = [
    "Railway云端部署异常排查（DNS可解析但443端口超时）",
    "本地代码推送GitHub，确保云端与本地代码同步",
    "前端操作层面的手工测试（API已通过，UI可能有小问题）",
    "等待王琛前辈试用反馈后进行针对性调整",
]
add_multi_text(slide, Inches(6.8), Inches(2.1), Inches(5.5), Inches(3.5), todo_items, font_size=12, prefix='⚠️ ', color=C_AMBER)

# ===========================================================
# Slide: 结束页
# ===========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_PRIMARY)
add_textbox(slide, Inches(1), Inches(2.5), Inches(10), Inches(0.8),
            "感谢聆听", font_size=32, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.67), Inches(3.5), Inches(2), Inches(0.04), C_WHITE)
add_textbox(slide, Inches(1), Inches(3.8), Inches(10), Inches(0.4),
            "入职19天 · 从零到一构建MRP II系统", font_size=16, color=C_WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.4), Inches(10), Inches(0.4),
            "在实战中学习，在迭代中成长", font_size=14, color=C_WHITE, alignment=PP_ALIGN.CENTER)

# ===========================================================
# 保存
# ===========================================================
output_path = r"C:\Users\Admin\Desktop\金崧的日报\MRP系统周工作总结-20260701-0707.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
